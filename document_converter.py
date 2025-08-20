import os
import tempfile
import logging
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import black
from reportlab.lib.units import inch

logger = logging.getLogger(__name__)

def convert_document_to_pdf(file_path, file_name):
    """Convert various document formats to PDF"""
    
    file_extension = os.path.splitext(file_name)[1].lower()
    
    if file_extension == '.docx':
        return convert_docx_to_pdf(file_path)
    elif file_extension == '.xlsx':
        return convert_xlsx_to_pdf(file_path)
    elif file_extension == '.pptx':
        return convert_pptx_to_pdf(file_path)
    elif file_extension == '.html':
        return convert_html_to_pdf(file_path)
    elif file_extension == '.txt':
        return convert_txt_to_pdf(file_path)
    else:
        raise ValueError(f"Unsupported file format: {file_extension}")

def convert_docx_to_pdf(docx_path):
    """Convert Word document to PDF"""
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pdf')
    os.close(temp_fd)
    
    try:
        from docx import Document
        
        # Read Word document
        doc = Document(docx_path)
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        normal_style = styles['Normal']
        heading_style = styles['Heading1']
        
        story = []
        
        # Process paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Determine if it's a heading (bold and larger text)
                is_heading = False
                if paragraph.runs:
                    for run in paragraph.runs:
                        if run.bold and len(paragraph.text.strip()) < 100:
                            is_heading = True
                            break
                
                # Create paragraph
                if is_heading:
                    para = Paragraph(paragraph.text.strip(), heading_style)
                else:
                    para = Paragraph(paragraph.text.strip(), normal_style)
                
                story.append(para)
                story.append(Spacer(1, 6))
        
        # Process tables
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                table_data.append(row_data)
            
            if table_data:
                # Create table
                pdf_table = Table(table_data)
                pdf_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), '#CCCCCC'),
                    ('TEXTCOLOR', (0, 0), (-1, 0), black),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), '#F7F7F7'),
                    ('GRID', (0, 0), (-1, -1), 1, black)
                ]))
                story.append(pdf_table)
                story.append(Spacer(1, 12))
        
        # Build PDF
        pdf_doc.build(story)
        
        logger.info(f"DOCX converted to PDF: {temp_path}")
        return temp_path
        
    except Exception as e:
        logger.error(f"Error converting DOCX to PDF: {e}")
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise

def convert_xlsx_to_pdf(xlsx_path):
    """Convert Excel spreadsheet to PDF"""
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pdf')
    os.close(temp_fd)
    
    try:
        from openpyxl import load_workbook
        
        # Load workbook
        workbook = load_workbook(xlsx_path)
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            rightMargin=36,
            leftMargin=36,
            topMargin=72,
            bottomMargin=18
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        heading_style = styles['Heading2']
        
        story = []
        
        # Process each worksheet
        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            
            # Add sheet title
            if len(workbook.sheetnames) > 1:
                title = Paragraph(f"Sheet: {sheet_name}", heading_style)
                story.append(title)
                story.append(Spacer(1, 12))
            
            # Get data range
            max_row = worksheet.max_row
            max_col = worksheet.max_column
            
            if max_row > 0 and max_col > 0:
                # Extract data
                table_data = []
                for row in range(1, min(max_row + 1, 101)):  # Limit to 100 rows
                    row_data = []
                    for col in range(1, min(max_col + 1, 11)):  # Limit to 10 columns
                        cell = worksheet.cell(row=row, column=col)
                        value = str(cell.value) if cell.value is not None else ""
                        row_data.append(value[:50])  # Limit cell content length
                    table_data.append(row_data)
                
                # Create table
                if table_data:
                    pdf_table = Table(table_data)
                    pdf_table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), '#CCCCCC'),
                        ('TEXTCOLOR', (0, 0), (-1, 0), black),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('BOTTOMPADDING', (0, 0), (-1, 0), 8),
                        ('BACKGROUND', (0, 1), (-1, -1), '#F7F7F7'),
                        ('GRID', (0, 0), (-1, -1), 0.5, black)
                    ]))
                    story.append(pdf_table)
                    story.append(Spacer(1, 20))
        
        # Build PDF
        pdf_doc.build(story)
        
        logger.info(f"XLSX converted to PDF: {temp_path}")
        return temp_path
        
    except Exception as e:
        logger.error(f"Error converting XLSX to PDF: {e}")
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise

def convert_pptx_to_pdf(pptx_path):
    """Convert PowerPoint presentation to PDF"""
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pdf')
    os.close(temp_fd)
    
    try:
        from pptx import Presentation
        
        # Load presentation
        presentation = Presentation(pptx_path)
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = styles['Title']
        heading_style = styles['Heading2']
        normal_style = styles['Normal']
        
        story = []
        
        # Process each slide
        for slide_num, slide in enumerate(presentation.slides, 1):
            # Add slide title
            slide_title = Paragraph(f"Slide {slide_num}", title_style)
            story.append(slide_title)
            story.append(Spacer(1, 12))
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text = shape.text_frame.text.strip()
                elif hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                else:
                    continue
                    
                if text:
                    # Determine style based on text characteristics
                    if len(text) < 100 and slide_num == 1:
                        # Likely a title
                        para = Paragraph(text, heading_style)
                    else:
                        para = Paragraph(text, normal_style)
                    
                    story.append(para)
                    story.append(Spacer(1, 6))
            
            # Add separator between slides
            if slide_num < len(presentation.slides):
                story.append(Spacer(1, 20))
        
        # Build PDF
        pdf_doc.build(story)
        
        logger.info(f"PPTX converted to PDF: {temp_path}")
        return temp_path
        
    except Exception as e:
        logger.error(f"Error converting PPTX to PDF: {e}")
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise

def convert_html_to_pdf(html_path):
    """Convert HTML file to PDF"""
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pdf')
    os.close(temp_fd)
    
    try:
        from bs4 import BeautifulSoup
        
        # Read HTML file
        with open(html_path, 'r', encoding='utf-8') as file:
            html_content = file.read()
        
        # Parse HTML
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        title_style = styles['Title']
        heading_style = styles['Heading1']
        normal_style = styles['Normal']
        
        story = []
        
        # Extract title
        title = soup.find('title')
        if title and title.get_text().strip():
            title_para = Paragraph(title.get_text().strip(), title_style)
            story.append(title_para)
            story.append(Spacer(1, 20))
        
        # Process body content
        body = soup.find('body') or soup
        
        for element in body.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'span']):
            if element and hasattr(element, 'get_text'):
                text = element.get_text().strip()
                if text:
                    # Determine style based on tag
                    if hasattr(element, 'name') and element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                        para = Paragraph(text, heading_style)
                    else:
                        para = Paragraph(text, normal_style)
                    
                    story.append(para)
                    story.append(Spacer(1, 6))
        
        # Build PDF
        pdf_doc.build(story)
        
        logger.info(f"HTML converted to PDF: {temp_path}")
        return temp_path
        
    except Exception as e:
        logger.error(f"Error converting HTML to PDF: {e}")
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise

def convert_txt_to_pdf(txt_path):
    """Convert text file to PDF"""
    
    # Create temporary file
    temp_fd, temp_path = tempfile.mkstemp(suffix='.pdf')
    os.close(temp_fd)
    
    try:
        # Read text file
        with open(txt_path, 'r', encoding='utf-8') as file:
            text_content = file.read()
        
        # Create PDF document
        pdf_doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            rightMargin=72,
            leftMargin=72,
            topMargin=72,
            bottomMargin=18
        )
        
        # Get styles
        styles = getSampleStyleSheet()
        normal_style = styles['Normal']
        
        story = []
        
        # Split text into paragraphs
        paragraphs = text_content.split('\n')
        
        for para_text in paragraphs:
            if para_text.strip():
                para = Paragraph(para_text.strip(), normal_style)
                story.append(para)
                story.append(Spacer(1, 6))
        
        # Build PDF
        pdf_doc.build(story)
        
        logger.info(f"TXT converted to PDF: {temp_path}")
        return temp_path
        
    except Exception as e:
        logger.error(f"Error converting TXT to PDF: {e}")
        if os.path.exists(temp_path):
            os.unlink(temp_path)
        raise
