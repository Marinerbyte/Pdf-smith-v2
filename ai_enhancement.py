import os
import logging
import json
import tempfile
from PyPDF2 import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
import cv2
import numpy as np

# Try to import Groq - handle import error gracefully
try:
    from groq import Groq
    GROQ_IMPORT_SUCCESS = True
except ImportError:
    GROQ_IMPORT_SUCCESS = False
    logger = logging.getLogger(__name__)
    logger.warning("Groq package not installed - AI features will be disabled")

logger = logging.getLogger(__name__)

# Initialize Groq client
client = None
try:
    GROQ_API_KEY = os.environ.get("GROQ_API_KEY")
    if GROQ_API_KEY and GROQ_IMPORT_SUCCESS:
        client = Groq(api_key=GROQ_API_KEY)
        AI_AVAILABLE = True
        logger.info("Groq AI client initialized successfully")
    else:
        AI_AVAILABLE = False
        if not GROQ_IMPORT_SUCCESS:
            logger.warning("Groq package not available - AI features disabled")
        else:
            logger.warning("GROQ_API_KEY not found - AI features disabled")
except Exception as e:
    logger.error(f"Error initializing Groq client: {e}")
    AI_AVAILABLE = False

def extract_text_from_pdf(file_path):
    """Extract text content from PDF file"""
    try:
        reader = PdfReader(file_path)
        text_content = []
        
        for page_num, page in enumerate(reader.pages, 1):
            text = page.extract_text()
            if text.strip():
                text_content.append(f"Page {page_num}:\n{text.strip()}")
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from PDF: {e}")
        return None

def extract_text_from_docx(file_path):
    """Extract text content from Word document"""
    try:
        doc = Document(file_path)
        text_content = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX: {e}")
        return None

def extract_text_from_xlsx(file_path):
    """Extract text content from Excel file"""
    try:
        workbook = load_workbook(file_path)
        text_content = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = f"Sheet: {sheet_name}\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:
                    sheet_text += " | ".join(row_text) + "\n"
            
            text_content.append(sheet_text)
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from XLSX: {e}")
        return None

def extract_text_from_pptx(file_path):
    """Extract text content from PowerPoint presentation"""
    try:
        prs = Presentation(file_path)
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = f"Slide {slide_num}:\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and hasattr(shape.text, 'strip') and shape.text.strip():
                    slide_text += str(shape.text).strip() + "\n"
            
            if slide_text.strip() != f"Slide {slide_num}:":
                text_content.append(slide_text)
        
        return "\n\n".join(text_content)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX: {e}")
        return None

def extract_text_from_image(file_path):
    """Extract text from image using OCR"""
    try:
        # Read image
        img = cv2.imread(file_path)
        if img is None:
            return None
        
        # Convert to grayscale and apply threshold
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Extract text using pytesseract
        text = pytesseract.image_to_string(thresh, config='--psm 6')
        
        return text.strip() if text.strip() else None
    except Exception as e:
        logger.error(f"Error extracting text from image: {e}")
        return None

def analyze_document_with_ai(content, document_type="document"):
    """Analyze document content using Groq AI and provide enhancement suggestions"""
    if not AI_AVAILABLE:
        return {
            "error": "AI analysis not available. Please set GROQ_API_KEY environment variable.",
            "suggestions": []
        }
    
    if not content or len(content.strip()) < 50:
        return {
            "error": "Document content is too short for meaningful analysis.",
            "suggestions": ["Add more detailed content to get better enhancement suggestions."]
        }
    
    # Truncate content if too long (Groq has token limits)
    max_content_length = 8000  # Conservative limit for Llama 7B
    if len(content) > max_content_length:
        content = content[:max_content_length] + "\n...[Content truncated for analysis]"
    
    try:
        prompt = f"""
You are an expert document enhancement specialist. Please provide an enhanced, improved version of this {document_type} document.

ORIGINAL DOCUMENT:
{content}

Please create an enhanced version that:
- Has better structure and organization
- Uses clearer, more professional language
- Improves formatting and readability
- Adds relevant headings and sections
- Enhances overall presentation

Provide your response in this format:

**ENHANCED DOCUMENT:**
[Write the complete improved version of the document here - make it significantly better with proper structure, clear language, and professional presentation]

**IMPROVEMENTS MADE:**
- [List the key improvements you made]
- [Keep this list concise - 3-5 main points]

Focus on creating a polished, professional version that the user can immediately use.
"""
        
        # Use Groq API with Llama 7B
        if not client:
            raise Exception("Groq client not initialized")
        
        chat_completion = client.chat.completions.create(
            messages=[
                {
                    "role": "user",
                    "content": prompt,
                }
            ],
            model="llama3-8b-8192",  # Using Llama 3 8B for better performance
            temperature=0.7,
            max_tokens=2048,
            top_p=0.9
        )
        
        response_text = chat_completion.choices[0].message.content.strip()
        
        # Return the enhanced document directly
        return {
            "success": True,
            "enhanced_document": response_text,
            "ai_model": "Llama 3 8B (via Groq)"
        }
            
    except Exception as e:
        logger.error(f"Error in AI analysis: {e}")
        return {
            "error": f"AI analysis failed: {str(e)}",
            "suggestions": ["Please try again or check your document format."]
        }

def format_enhancement_suggestions(analysis_result):
    """Format AI analysis results into user-friendly text"""
    if "error" in analysis_result:
        return f"❌ **Analysis Error**\n{analysis_result['error']}"
    
    if not analysis_result.get("success"):
        return "❌ **Analysis Failed**\nUnable to analyze document content."
    
    analysis = analysis_result.get("analysis", {})
    ai_model = analysis_result.get("ai_model", "AI Assistant")
    
    # Handle enhanced document format
    if "enhanced_document" in analysis_result:
        return f"🤖 **AI Document Enhancement** (powered by {ai_model})\n\n{analysis_result['enhanced_document']}"
    
    # Format structured analysis
    formatted_text = f"🤖 **AI Document Enhancement Analysis**\n*Powered by {ai_model}*\n\n"
    
    # Overall score
    if "overall_score" in analysis:
        score = analysis["overall_score"]
        formatted_text += f"📊 **Overall Score:** {score}/10\n\n"
    
    # Summary
    if "summary" in analysis:
        formatted_text += f"📋 **Summary:**\n{analysis['summary']}\n\n"
    
    # Categories
    categories = {
        "structure": "🏗️ **Structure & Organization:**",
        "clarity": "✨ **Clarity & Readability:**",
        "content": "📝 **Content Enhancement:**",
        "formatting": "🎨 **Formatting & Visual Appeal:**",
        "professionalism": "💼 **Professional Polish:**"
    }
    
    for category, title in categories.items():
        if category in analysis and analysis[category]:
            formatted_text += f"{title}\n"
            for suggestion in analysis[category]:
                formatted_text += f"• {suggestion}\n"
            formatted_text += "\n"
    
    return formatted_text

def analyze_document_file(file_path, file_type):
    """Main function to analyze any supported document type"""
    try:
        # Extract text based on file type
        if file_type.lower() in ['pdf']:
            content = extract_text_from_pdf(file_path)
            doc_type = "PDF document"
        elif file_type.lower() in ['docx', 'doc']:
            content = extract_text_from_docx(file_path)
            doc_type = "Word document"
        elif file_type.lower() in ['xlsx', 'xls']:
            content = extract_text_from_xlsx(file_path)
            doc_type = "Excel spreadsheet"
        elif file_type.lower() in ['pptx', 'ppt']:
            content = extract_text_from_pptx(file_path)
            doc_type = "PowerPoint presentation"
        elif file_type.lower() in ['jpg', 'jpeg', 'png', 'gif', 'webp']:
            content = extract_text_from_image(file_path)
            doc_type = "image with text"
        else:
            # Try to read as plain text
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            doc_type = "text document"
        
        if not content:
            return {
                "error": f"Could not extract readable content from the {doc_type}.",
                "suggestions": ["Make sure the file contains text content or is not corrupted."]
            }
        
        # Analyze with AI
        return analyze_document_with_ai(content, doc_type)
        
    except Exception as e:
        logger.error(f"Error analyzing document: {e}")
        return {
            "error": f"Error processing document: {str(e)}",
            "suggestions": ["Please try uploading a different file or check if the file is corrupted."]
        }